#!/usr/bin/env python
"""Fill the IDBI Innovate Prototype Submission Deck with the IDBI Wealth Copilot content."""

import os
from PIL import Image, ImageChops
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = "/Users/rahulsingh/Desktop/IDBI Wealth Copilot/docs"
DECK = os.path.join(ROOT, "Prototype Submission Deck _ IDBI Innovate.pptx")
ASSETS = os.path.join(ROOT, "deck_assets")
SHOTS = os.path.join(ROOT, "screenshots")

NAVY = RGBColor(0x1E, 0x27, 0x61)
TEAL = RGBColor(0x0B, 0x7A, 0x70)
ACCENT = RGBColor(0x0B, 0x8F, 0x86)
INK = RGBColor(0x24, 0x2C, 0x3A)
MUTED = RGBColor(0x5A, 0x66, 0x7C)
LINKC = RGBColor(0x15, 0x5E, 0x9A)
FONT = "Calibri"


# ---- crop helper for the wide process-flow png ----
def crop_to_content(path):
    src = Image.open(path).convert("RGBA")
    flat = Image.new("RGBA", src.size, (255, 255, 255, 255))
    flat.alpha_composite(src)
    im = flat.convert("RGB")
    white = Image.new("RGB", im.size, (255, 255, 255))
    diff = ImageChops.difference(im, white).convert("L").point(lambda v: 255 if v > 20 else 0)
    bbox = diff.getbbox()
    if bbox:
        x0, y0, x1, y1 = bbox
        p = 6
        im = im.crop((max(0, x0 - p), max(0, y0 - p), min(im.size[0], x1 + p), min(im.size[1], y1 + p)))
    im.save(path, "PNG")
    return im.size


crop_to_content(os.path.join(ASSETS, "process-flow.png"))
crop_to_content(os.path.join(ASSETS, "architecture.png"))


# ---- text helpers ----
def set_para(p, runs, align=PP_ALIGN.LEFT, space_after=5, space_before=0):
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    for text, size, bold, color in runs:
        r = p.add_run()
        r.text = text
        f = r.font
        f.size = Pt(size)
        f.bold = bold
        f.color.rgb = color
        f.name = FONT


def add_box(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Inches(0.04))
    return tf


def add_paras(tf, items):
    first = True
    for runs in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        set_para(p, runs)


def header(text):
    return [(text, 13, True, ACCENT)]


def body(text, size=11.5):
    return [(text, size, False, INK)]


def lead(label, rest, size=11.5):
    return [(label, size, True, NAVY), (rest, size, False, INK)]


def place_image(slide, path, top, max_w, max_h, cx=5.0):
    iw, ih = Image.open(path).size
    r = iw / ih
    w = max_w
    h = w / r
    if h > max_h:
        h = max_h
        w = h * r
    left = cx - w / 2
    slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(w), Inches(h))
    return left, top, w, h


def caption(slide, text, l, t, w):
    tf = add_box(slide, l, t, w, 0.3)
    set_para(tf.paragraphs[0], [(text, 9, False, MUTED)], align=PP_ALIGN.CENTER, space_after=0)


prs = Presentation(DECK)
S = list(prs.slides)

# ===== SLIDE 1 — Team details (rebuild the existing text box) =====
box = S[0].shapes[1]
box.height = Inches(1.95)
tf = box.text_frame
tf.clear()
add_paras(tf, [
    [("Team Details", 16, True, NAVY)],
    [("Team name:  ", 13, True, INK), ("Agentic_Hackers", 13, False, INK)],
    [("Team leader:  ", 13, True, INK), ("Gokul D", 13, False, INK)],
    [("Team size:  ", 13, True, INK), ("5", 13, False, INK)],
    [("Problem Statement — Track 01 · Wealth Advisory", 12.5, True, ACCENT)],
    body("Wealth management and advisory remain fragmented and inaccessible to most customers. "
         "The absence of comprehensive insight into customer investment behaviour and spending "
         "limits timely, personalized, data-driven guidance.", 11),
])

# ===== SLIDE 2 — Brief about the idea =====
tf = add_box(S[1], 0.4, 1.55, 9.2, 3.8)
add_paras(tf, [
    header("What it is"),
    body("IDBI Wealth Copilot (“Aanya”) is an avatar-led, agentic AI wealth advisor embedded in "
         "the IDBI mobile app. It turns each customer’s own transaction behaviour into personalized, "
         "explainable, bank-grade guidance — through a conversational 3D avatar with voice and "
         "visual, chart-first answers."),
    [("", 6, False, INK)],
    header("Why it’s different"),
    body("Unlike a chatbot, Aanya runs on a real quantitative engine (risk profiling, goal-based SIP "
         "planning, portfolio construction, behavioural insights) wrapped in compliance guardrails "
         "enforced in code: she never promises returns and only recommends products she has actually "
         "looked up."),
    [("", 6, False, INK)],
    header("The outcome"),
    [("A personal wealth manager for every customer — at near-zero marginal cost.", 12, True, NAVY)],
])

# ===== SLIDE 3 — Opportunities =====
tf = add_box(S[2], 0.4, 2.02, 9.2, 3.4)
add_paras(tf, [
    [("How different is it?", 12.5, True, ACCENT)],
    body("Most rivals are generic robo-advisors or LLM chatbots. Aanya combines three things they "
         "don’t: a real, auditable quant engine (every number computed, not generated); compliance "
         "guardrails enforced in code and unit-tested; and an avatar + voice + chart-first UX built "
         "for the IDBI app.", 11),
    [("How will it solve the problem?", 12.5, True, ACCENT)],
    body("It reads transactions to estimate income, savings rate, idle cash and costly debt, then "
         "plans goals and recommends only real IDBI catalogue products — making personalized "
         "advisory available to every customer, not just the affluent.", 11),
    [("USP of the solution", 12.5, True, ACCENT)],
    body("Explainable + compliant by design: it shows the “why” behind every recommendation, "
         "refuses to guarantee returns, and is one config switch away from IDBI’s sandbox APIs.", 11),
])

# ===== SLIDE 4 — Features =====
tf = add_box(S[3], 0.4, 1.55, 9.2, 3.85)
add_paras(tf, [
    lead("Conversational 3D avatar — ", "voice in & out, voice-reactive, multilingual-ready."),
    lead("Behavioural insights — ", "income, savings-rate, idle-cash and costly-debt detection with ranked next-best-actions."),
    lead("Goal-based planning — ", "required SIP, asset allocation and projections with a live what-if slider."),
    lead("Risk profiling — ", "an explainable 0–100 score that drives the recommended allocation."),
    lead("Grounded recommendations — ", "products retrieved from the IDBI catalogue (RAG); indicative returns only."),
    lead("Chart-first answers — ", "donut, line, bar and gauge visuals for instant understanding."),
    lead("Compliance guardrails — ", "blocks guaranteed-return claims and ungrounded products (code-enforced, unit-tested)."),
    lead("Bank-app UX — ", "light/dark, responsive; advice only, never executes transactions."),
])

# ===== SLIDE 5 — Process flow (image) =====
place_image(S[4], os.path.join(ASSETS, "process-flow.png"), top=1.7, max_w=9.0, max_h=3.35)
caption(S[4], "Data questions are answered instantly as charts (offline); open-ended questions use the guarded LLM agent.",
        0.4, 5.12, 9.2)

# ===== SLIDE 6 — Wireframe / mock (image) =====
l, t, w, h = place_image(S[5], os.path.join(SHOTS, "02-goal-plan.png"), top=1.55, max_w=4.6, max_h=3.6)
caption(S[5], "Goal-planning screen — allocation donut, projected-growth chart and a live what-if slider (mobile-style UI).",
        0.4, 5.18, 9.2)

# ===== SLIDE 7 — Architecture (image) =====
place_image(S[6], os.path.join(ASSETS, "architecture.png"), top=1.5, max_w=6.6, max_h=3.7)
caption(S[6], "Four layers; everything above the provider seam is independent of the data source (synthetic now → IDBI sandbox later).",
        0.4, 5.28, 9.2)

# ===== SLIDE 8 — Technologies =====
tf = add_box(S[7], 0.4, 1.55, 9.2, 3.85)
add_paras(tf, [
    lead("Frontend — ", "HTML/CSS/JS, Three.js (3D avatar), Chart.js (visuals), Web Speech API (voice)."),
    lead("Backend — ", "Python 3, FastAPI, Pydantic v2."),
    lead("AI — ", "Anthropic Claude (claude-opus-4-8) tool-calling agent with adaptive thinking."),
    lead("Quant / ML — ", "pure-Python engine (risk, goals, portfolio, insights); TF-IDF product retrieval (RAG)."),
    lead("Data — ", "seeded synthetic datasets; provider seam ready for IDBI sandbox APIs."),
    lead("Cloud — ", "AWS / ACC-ready, containerizable, stateless API."),
    lead("Quality — ", "20 automated tests (pytest); deterministic, reproducible builds."),
])

# ===== SLIDE 9 — Estimated implementation cost =====
tf = add_box(S[8], 0.4, 1.55, 9.2, 3.85)
add_paras(tf, [
    body("Round-1 prototype runs at effectively zero infrastructure cost (synthetic data, "
         "browser-native voice and avatar). Indicative production economics:"),
    [("", 5, False, INK)],
    lead("LLM inference (Claude) — ", "usage-based; prompt caching and a Sonnet fallback keep cost low; only open-ended chats hit the LLM."),
    lead("Cloud hosting (AWS / ACC) — ", "covered by hackathon credits during the pilot; scales with traffic."),
    lead("Voice & avatar — ", "browser-native Web Speech = ₹0; an optional premium avatar (D-ID / HeyGen) is usage-based."),
    lead("Data & quant engine — ", "runs in-process; no per-query cost."),
    [("", 5, False, INK)],
    [("Net: marginal cost per advised customer approaches zero versus a human relationship manager.", 12, True, NAVY)],
])

# ===== SLIDE 10 — Snapshots (three screenshots) =====
shots = [
    ("01-snapshot-charts.png", "Snapshot — spending & actions"),
    ("02-goal-plan.png", "Goal plan — allocation & growth"),
    ("03-products.png", "Products — indicative returns"),
]
cols = [0.4, 3.5, 6.6]
colw = 3.0
for (fname, cap), cl in zip(shots, cols):
    path = os.path.join(SHOTS, fname)
    iw, ih = Image.open(path).size
    w = colw
    h = w * ih / iw
    if h > 3.4:
        h = 3.4
        w = h * iw / ih
    left = cl + (colw - w) / 2
    S[9].shapes.add_picture(path, Inches(left), Inches(1.6), Inches(w), Inches(h))
    caption(S[9], cap, cl, 1.6 + h + 0.04, colw)

# ===== SLIDE 11 — Performance / benchmarking =====
tf = add_box(S[10], 0.4, 1.55, 9.2, 3.85)
add_paras(tf, [
    lead("API latency — ", "quant endpoints respond in <50 ms (in-memory, pure-Python) — instant chart answers."),
    lead("Advisory turn — ", "the LLM + tool loop completes in a few seconds with adaptive thinking."),
    lead("Test coverage — ", "20/20 automated tests passing (quant math, guardrails, API)."),
    lead("Guardrail efficacy — ", "100% block rate on guaranteed-return claims in the test suite."),
    lead("Advisory reach — ", "from a handful of RM-served clients to every customer — advisory is no longer human-limited."),
    lead("Reproducibility — ", "seeded synthetic data yields identical results across runs."),
])

# ===== SLIDE 12 — Future development =====
tf = add_box(S[11], 0.4, 1.55, 9.2, 3.85)
add_paras(tf, [
    lead("Live integration — ", "flip the provider seam to IDBI sandbox APIs for real accounts and transactions."),
    lead("Richer avatar — ", "premium photoreal avatar and full multilingual voice (Hindi + regional)."),
    lead("Market data — ", "real-time NAVs and an expanded product catalogue."),
    lead("Proactive advice — ", "portfolio monitoring, rebalancing alerts and next-best-action nudges."),
    lead("Cross-sell — ", "income-estimation and creditworthiness models (Track 02 synergy)."),
    lead("Compliance — ", "RBI/SEBI hardening, audit logging and on-device privacy."),
    lead("RM co-pilot — ", "an advisor-assist mode and A/B-tested behavioural nudges."),
])

# ===== SLIDE 13 — Links (append to existing labels) =====
links_box = S[12].shapes[0].text_frame
for p in links_box.paragraphs:
    t = "".join(r.text for r in p.runs).strip()
    if t.startswith("GitHub"):
        r = p.add_run(); r.text = "   →   https://github.com/Rahulcse79/IDBI-Wealth-Copilot"
        r.font.size = Pt(13); r.font.bold = False; r.font.color.rgb = LINKC; r.font.name = FONT
    elif t.startswith("Demo Video"):
        r = p.add_run(); r.text = "   →   < add your 3-minute demo video link >"
        r.font.size = Pt(13); r.font.italic = True; r.font.color.rgb = MUTED; r.font.name = FONT
    elif t.startswith("Final Product"):
        r = p.add_run(); r.text = "   →   < add your deployed product link >"
        r.font.size = Pt(13); r.font.italic = True; r.font.color.rgb = MUTED; r.font.name = FONT

prs.save(DECK)
print("Saved:", DECK)
print("Slides:", len(prs.slides.__iter__.__self__._sldIdLst))
